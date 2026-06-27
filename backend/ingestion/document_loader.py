import os
import logging
import pandas as pd
import pytesseract
from PIL import Image
from pdf2image import convert_from_path
from docx import Document as DocxDocument
from pptx import Presentation
from langchain_core.documents import Document
from langchain_community.document_loaders import PyPDFLoader

logger = logging.getLogger(__name__)

def load_document(file_path):
    filename = os.path.basename(file_path)

    extension = os.path.splitext(
        file_path
    )[1].lower()

    # PDF
    if extension == ".pdf":

        loader = PyPDFLoader(
            file_path
        )
        
        docs = loader.load()
        
        # Check if it's a scanned PDF (no text found)
        total_text = "".join([d.page_content for d in docs]).strip()
        if not total_text:
            logger.info("[Loader] PDF has no extracted text; falling back to OCR processing.")
            try:
                images = convert_from_path(file_path)
                ocr_docs = []
                for i, img in enumerate(images):
                    text = pytesseract.image_to_string(img)
                    ocr_docs.append(Document(
                        page_content=text,
                        metadata={"source": filename, "page": i + 1, "ocr_used": True}
                    ))
                logger.info("[Loader] Successfully completed OCR for scanned PDF: %s", filename)
                return ocr_docs
            except Exception as e:
                logger.error("[Loader] PDF OCR failed. System binaries (poppler or tesseract) may be missing.", exc_info=True)
                raise ValueError(
                    "This PDF appears to be scanned and requires server-side OCR. "
                    f"OCR processing failed because system packages are missing or misconfigured: {str(e)}"
                )

        return docs

    # DOCX
    elif extension == ".docx":

        doc = DocxDocument(
            file_path
        )

        text = "\n".join(
            [p.text for p in doc.paragraphs]
        )

        return [
            Document(
                page_content=text,
                metadata={
                    "source": filename
                }
            )
        ]

    # PPTX
    elif extension == ".pptx":

        prs = Presentation(
            file_path
        )

        slides_text = []

        for slide in prs.slides:

            for shape in slide.shapes:

                if hasattr(
                    shape,
                    "text"
                ):
                    slides_text.append(
                        shape.text
                    )

        return [
            Document(
                page_content="\n".join(
                    slides_text
                ),
                metadata={
                    "source": filename
                }
            )
        ]

    # TXT
    elif extension == ".txt":

        with open(
            file_path,
            "r",
            encoding="utf-8"
        ) as f:

            text = f.read()

        return [
            Document(
                page_content=text,
                metadata={
                    "source": filename
                }
            )
        ]

    # MD
    elif extension == ".md":

        with open(
            file_path,
            "r",
            encoding="utf-8"
        ) as f:

            text = f.read()

        return [
            Document(
                page_content=text,
                metadata={
                    "source": filename
                }
            )
        ]

    # CSV
    elif extension == ".csv":

        df = pd.read_csv(
            file_path
        )

        return [
            Document(
                page_content=df.to_string(),
                metadata={
                    "source": filename
                }
            )
        ]

    # XLSX
    elif extension == ".xlsx":

        df = pd.read_excel(
            file_path
        )

        return [
            Document(
                page_content=df.to_string(),
                metadata={
                    "source": filename
                }
            )
        ]

    # Images (OCR)
    elif extension in [".jpg", ".jpeg", ".png"]:
        logger.info("[Loader] Attempting OCR for image file: %s", filename)
        try:
            img = Image.open(file_path)
            text = pytesseract.image_to_string(img)
            logger.info("[Loader] Successfully completed OCR for image: %s", filename)
            return [
                Document(
                    page_content=text,
                    metadata={
                        "source": filename,
                        "ocr_used": True
                    }
                )
            ]
        except Exception as e:
            logger.error("[Loader] Image OCR failed. Tesseract binary may be missing.", exc_info=True)
            raise ValueError(
                "Image uploads require server-side OCR. "
                f"OCR processing failed because Tesseract is missing or misconfigured: {str(e)}"
            )

    else:

        raise ValueError(
            f"Unsupported file type: {extension}"
        )
